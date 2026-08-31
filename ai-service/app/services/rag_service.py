"""
RAG Service — Retrieves relevant material chunks for lesson grounding.

STUB IMPLEMENTATION: Returns an empty list until vector DB integration
is completed.  The lesson planner handles empty chunks gracefully by
falling back to topic-only generation.
"""

import logging

import os
import io
import logging
from typing import List, Dict, Any

# Text extraction
import pdfplumber
from docx import Document
from pptx import Presentation

# Vector DB
import chromadb
from chromadb.utils import embedding_functions

logger = logging.getLogger(__name__)

# Ensure ChromaDB path exists
CHROMA_DB_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "chroma_db")
os.makedirs(CHROMA_DB_PATH, exist_ok=True)

# Initialize ChromaDB client (persistent)
chroma_client = chromadb.PersistentClient(path=CHROMA_DB_PATH)

# Using default all-MiniLM-L6-v2 embedding function built into Chroma.
# For production, we can swap this for google-genai embeddings.
embedding_fn = embedding_functions.DefaultEmbeddingFunction()

# Create or get collection
collection = chroma_client.get_or_create_collection(
    name="ai_mentor_materials",
    embedding_function=embedding_fn
)

def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> List[str]:
    """Simple recursive character text splitter."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        if end < len(text):
            # Try to find a natural break point (newline or period)
            last_newline = text.rfind('\n', start, end)
            last_period = text.rfind('. ', start, end)
            break_point = max(last_newline, last_period)
            
            if break_point != -1 and break_point > start + (chunk_size // 2):
                end = break_point + 1 # Include the newline or period
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
            
        start = end - overlap
        if start < 0:
            start = 0
            
        # Prevent infinite loops if no progress is made
        if start >= end:
            start = end
            
    return chunks

async def process_material(file_bytes: bytes, filename: str, material_id: str):
    """
    Extract text from a file buffer, chunk it, and store in ChromaDB.
    Supports: .txt, .pdf, .docx, .pptx
    """
    logger.info("Processing material: %s (id: %s)", filename, material_id)
    text_content = ""
    ext = os.path.splitext(filename)[1].lower()

    try:
        if ext == '.txt':
            text_content = file_bytes.decode('utf-8')
            
        elif ext == '.pdf':
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                pages = [page.extract_text() or "" for page in pdf.pages]
                text_content = "\n\n".join(pages)
                
        elif ext == '.docx':
            doc = Document(io.BytesIO(file_bytes))
            text_content = "\n".join([para.text for para in doc.paragraphs])
            
        elif ext == '.pptx':
            prs = Presentation(io.BytesIO(file_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            text_content = "\n".join(text_runs)
            
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    except Exception as e:
        logger.error("Failed to extract text from %s: %s", filename, str(e))
        raise ValueError(f"Failed to extract text: {str(e)}")

    if not text_content.strip():
        logger.warning("No text extracted from %s", filename)
        return

    # Chunk the text
    chunks = chunk_text(text_content)
    logger.info("Extracted %d chunks from %s", len(chunks), filename)

    if not chunks:
        return

    # Prepare for ChromaDB
    ids = [f"{material_id}_{i}" for i in range(len(chunks))]
    metadatas = [{"material_id": material_id, "chunk_index": i, "source": filename} for i in range(len(chunks))]
    
    # Upsert into collection
    collection.upsert(
        documents=chunks,
        metadatas=metadatas,
        ids=ids
    )
    logger.info("Successfully ingested material %s into ChromaDB", material_id)


async def retrieve_chunks(material_id: str, query: str = None, top_k: int = 5) -> list[str]:
    """
    Retrieve the most relevant text chunks for a given material and query.
    If no query is provided, it tries to just return some chunks.
    """
    logger.info("Retrieving chunks for material_id=%s, query=%s", material_id, query)
    
    formatted_chunks = []
    
    # If a specific query is provided, we do similarity search
    if query:
        results = collection.query(
            query_texts=[query],
            n_results=top_k,
            where={"material_id": material_id}
        )
        if results and results["documents"] and results["documents"][0]:
            docs = results["documents"][0]
            metas = results["metadatas"][0]
            for doc, meta in zip(docs, metas):
                formatted_chunks.append(f"[Source: {meta.get('source', 'Unknown')}, Chunk Index: {meta.get('chunk_index', 0)}]\n{doc}")
            return formatted_chunks
        return []
    
    # Otherwise, just pull the first N chunks
    results = collection.get(
        where={"material_id": material_id},
        limit=top_k
    )
    
    if results and results["documents"]:
        docs = results["documents"]
        metas = results["metadatas"]
        for doc, meta in zip(docs, metas):
            formatted_chunks.append(f"[Source: {meta.get('source', 'Unknown')}, Chunk Index: {meta.get('chunk_index', 0)}]\n{doc}")
        return formatted_chunks
        
    return []
